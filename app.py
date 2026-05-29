#!/usr/bin/env python3
"""
Offline PDF OCR Web App with Auto-Language Detection
Flask application for rendering PDF OCR via web interface
"""

import os
import tempfile
import threading
import uuid
import re
import logging
import time
import shutil
import json
from flask import Flask, request, render_template, send_file, flash, redirect, url_for, jsonify
from werkzeug.utils import secure_filename
from pdf2image import convert_from_path, pdfinfo_from_path
import pytesseract
from pytesseract import image_to_osd
from PIL import Image
import requests
import gc

# Ensure pytesseract can find the tesseract binary installed via build.sh
for _tess in ['/usr/bin/tesseract', '/usr/local/bin/tesseract']:
    if os.path.isfile(_tess):
        pytesseract.pytesseract.tesseract_cmd = _tess
        break
try:
    from deep_translator import GoogleTranslator
    _DEEP_TRANSLATOR_AVAILABLE = True
except ImportError:
    _DEEP_TRANSLATOR_AVAILABLE = False
    logger.warning("deep-translator not installed; translation will be unavailable")

# Tell mega.py to use pycryptodome instead of the broken pycrypto
os.environ['MEGA_USE_CRYPTO_DOME'] = '1'
gc.set_threshold(100, 5, 2)  # More aggressive GC for memory-constrained environments
import asyncio
from concurrent.futures import ThreadPoolExecutor, TimeoutError as _CTimeoutError
if not hasattr(asyncio, 'coroutine'):
    asyncio.coroutine = lambda f: f
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", os.urandom(24).hex())

# Jinja filter: Unix timestamp to readable date
@app.template_filter('datetimeformat')
def datetimeformat(timestamp):
    import datetime
    return datetime.datetime.fromtimestamp(timestamp).strftime('%Y-%m-%d %H:%M')

# Configuration
MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE  # Reject oversized uploads early
ALLOWED_EXTENSIONS = {
    # PDF
    'pdf',
    # Images (OCR needed)
    'jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif',
    # Word documents
    'docx', 'doc',
    # PowerPoint
    'pptx', 'ppt',
    # Excel/Spreadsheets
    'xlsx', 'xls', 'csv',
    # Rich Text
    'rtf',
    # OpenDocument formats
    'odt', 'ods', 'odp',
    # HTML
    'html', 'htm',
    # XML/JSON
    'xml', 'json',
    # Plain text
    'txt', 'md'
}
DEFAULT_LANG = 'tam'  # Tamil by default
CHECKPOINT_INTERVAL = 5  # Upload checkpoint to Mega every N pages
BATCH_SIZE = 1  # Pages per batch (lower = less memory per batch)
MEGA_LOGIN_TIMEOUT = 30  # Seconds before Mega login times out
CONVERT_TIMEOUT = 120  # Max seconds for convert_from_path per batch
OCR_TIMEOUT = 300  # Max seconds per page for Tesseract OCR

# Progress tracking
progress_lock = threading.Lock()
progress_tracker = {}  # task_id: { ... }

# JSON file persistence for progress_tracker (survives Render restarts)
PROGRESS_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "progress_tracker.json")
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ocr-outputs")
_last_save_time = 0
_mega_restore_done = threading.Event()

def _save_progress(force=False):
    """Save progress_tracker to JSON file with throttling (max 1 write/sec)."""
    global _last_save_time
    now = time.time()
    if not force and now - _last_save_time < 2:
        return
    serializable = {}
    with progress_lock:
        for tid, task in progress_tracker.items():
            serializable[tid] = dict(task)
    try:
        with open(PROGRESS_FILE, 'w', encoding='utf-8') as f:
            json.dump(serializable, f, default=str, ensure_ascii=False)
        _last_save_time = now
    except Exception as e:
        logger.error(f"Failed to save progress: {e}")

def _load_progress():
    """Load progress_tracker from JSON file on startup."""
    if not os.path.exists(PROGRESS_FILE):
        return {}
    try:
        with open(PROGRESS_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        logger.error(f"Failed to load progress: {e}")
        return {}

def persist_output(task_id):
    """Copy completed output file from temp dir to persistent OUTPUT_DIR."""
    with progress_lock:
        task = progress_tracker.get(task_id)
        if not task:
            return False
        src = task.get("output_path")
        output_filename = task.get("output_filename")
        if not src or not os.path.exists(src) or not output_filename:
            return False
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    dest = os.path.join(OUTPUT_DIR, f"{task_id}_{output_filename}")
    try:
        shutil.copy2(src, dest)
        with progress_lock:
            task["output_path"] = dest
        logger.info(f"Task {task_id}: Output persisted to {dest}")
        return True
    except Exception as e:
        logger.error(f"Task {task_id}: Failed to persist output: {e}")
        return False


def rebuild_completed_from_local():
    """Scan local OUTPUT_DIR and restore completed tasks not in progress_tracker."""
    if not os.path.exists(OUTPUT_DIR):
        return
    now = time.time()
    restored = 0
    for fname in os.listdir(OUTPUT_DIR):
        if not fname.endswith(".txt"):
            continue
        filepath = os.path.join(OUTPUT_DIR, fname)
        if not os.path.isfile(filepath):
            continue
        # Filename format: {task_id}_{output_filename}
        task_id = fname.split("_", 1)[0]
        with progress_lock:
            if task_id in progress_tracker:
                continue
        output_filename = fname[len(task_id) + 1:]
        orig_name = output_filename.rsplit("_ocr.", 1)[0] if "_ocr." in output_filename else output_filename
        with progress_lock:
            progress_tracker[task_id] = {
                "current_page": 0, "status": "completed",
                "result_path": None, "error": None,
                "filename": orig_name, "output_filename": output_filename,
                "output_path": filepath,
                "download_link": None,
                "mega_uploaded": False, "mega_status": "",
                "file_type": "pdf", "detected_language": "",
                "pages_processed": 0, "percentage": 100,
                "download_count": 0, "completed_at": now, "created_at": now
            }
            restored += 1
    if restored:
        logger.info(f"Restored {restored} completed tasks from local {OUTPUT_DIR}")


def _get_memory_mb():
    """Get current RSS memory in MB (Linux only). Returns 0 on other platforms."""
    try:
        if os.name == 'posix':
            with open('/proc/self/status') as f:
                for line in f:
                    if line.startswith('VmRSS:'):
                        return int(line.split()[1]) // 1024
    except Exception:
        pass
    return 0

# Self-keepalive to prevent Render from sleeping during background tasks
_active_tasks = 0
_keepalive_lock = threading.Lock()
_keepalive_thread = None

# Cached Mega client (reused across calls to avoid repeated logins)
_mega_client = None
_mega_client_lock = threading.Lock()

def _get_mega_client():
    """Return cached Mega client, logging in only on first call."""
    global _mega_client
    if _mega_client is not None:
        return _mega_client
    email = os.environ.get("MEGA_EMAIL")
    password = os.environ.get("MEGA_PWD")
    if not email or not password:
        return None
    with _mega_client_lock:
        if _mega_client is not None:
            return _mega_client
        try:
            from mega import Mega
            _mega_client = Mega().login(email, password)
            logger.info("Mega client logged in (cached)")
        except Exception as e:
            logger.error(f"Mega login failed: {e}")
            _mega_client = None
    return _mega_client


def _translate_text(text, target_lang, source_lang='en'):
    """Translate text using Google Translate (via deep-translator, free, no API key). Returns translated string or None."""
    if not text or not text.strip() or not target_lang:
        return None
    if not _DEEP_TRANSLATOR_AVAILABLE:
        return None
    try:
        t = _get_translator(source_lang, target_lang)
        return t.translate(text)
    except Exception:
        logger.warning(f"Translation {source_lang}->{target_lang} failed", exc_info=True)
        return None


def _auto_translate_and_upload(task_id):
    """
    Upload pre-built translated output file to Mega ocr-translated/ folder.
    The translation is done page-by-page during OCR processing.
    Returns Mega download link or None.
    """
    with progress_lock:
        task = progress_tracker.get(task_id)
        if not task:
            return None
        target_lang = task.get("translate_to")
        translated_path = task.get("translated_output_path")
        base_filename = task.get("output_filename", "output.txt")
    if not target_lang or not translated_path or not os.path.exists(translated_path):
        return None
    try:
        translated_filename = base_filename.rsplit('.', 1)[0] + f"_{target_lang}.txt"
        mega_link = None
        m = _get_mega_client()
        if m:
            folder_handle = ensure_mega_folder(m, "ocr-translated")
            if folder_handle:
                mega_call(m, "upload", translated_path, dest=folder_handle, dest_filename=translated_filename, timeout=120)
                uploaded_node = mega_call(m, "find", f"ocr-translated/{translated_filename}")
                if uploaded_node:
                    mega_link = mega_call(m, "get_link", uploaded_node[0], timeout=30)
                    logger.info(f"Task {task_id}: Translated file uploaded to Mega: {mega_link}")
        return mega_link
    except Exception as e:
        logger.error(f"Task {task_id}: Auto-translate upload failed: {e}")
        return None


# Secondary keepalive: each active OCR thread also pings itself
_ocr_keepalive_running = threading.Event()

# Language detection mapping
SCRIPT_TO_LANG = {
    "Tamil": "tam",
    "Latin": "eng",  # Default Latin script to English
    "Devanagari": "hin",  # Hindi (most common Devanagari)
    "Telugu": "tel",
    "Bengali": "ben",
    "Kannada": "kan",
    "Malayalam": "mal",
    "Gujarati": "guj",
    "Punjabi": "pan",
    "Marathi": "mar",
    "Arabic": "ara",
    "Cyrillic": "rus",  # Russian
    "Greek": "ell",
    "Hebrew": "heb",
    "Thai": "tha",
    "Chinese": "chi_sim",  # Simplified Chinese
    "Japanese": "jpn",
    "Korean": "kor",
    "Spanish": "spa",
    "French": "fra",
    "German": "deu",
    "Italian": "ita",
    # Add more as needed
}

# Language code to full name mapping
LANG_CODE_TO_NAME = {
    "tam": "Tamil",
    "eng": "English",
    "hin": "Hindi",
    "tel": "Telugu",
    "ben": "Bengali",
    "kan": "Kannada",
    "mal": "Malayalam",
    "guj": "Gujarati",
    "pan": "Punjabi",
    "mar": "Marathi",
    "ara": "Arabic",
    "rus": "Russian",
    "ell": "Greek",
    "heb": "Hebrew",
    "tha": "Thai",
    "chi_sim": "Chinese (Simplified)",
    "jpn": "Japanese",
    "kor": "Korean",
    "spa": "Spanish",
    "fra": "French",
    "deu": "German",
    "ita": "Italian",
}

# Tesseract language code to ISO 639-1 code mapping (for translation source detection)
TESSERACT_TO_ISO = {
    "tam": "ta", "eng": "en", "hin": "hi", "mal": "ml",
    "tel": "te", "kan": "kn", "ben": "bn", "guj": "gu",
    "pan": "pa", "mar": "mr", "ara": "ar", "spa": "es",
    "fra": "fr", "deu": "de", "ita": "it", "rus": "ru",
    "chi_sim": "zh", "jpn": "ja", "kor": "ko", "tha": "th",
    "ell": "el", "heb": "he",
}

# Cache translator instance to avoid re-initializing
_translator_cache = {}

def _get_translator(source, target):
    """Get or create a cached GoogleTranslator instance."""
    key = f"{source}-{target}"
    if key not in _translator_cache:
        _translator_cache[key] = GoogleTranslator(source=source, target=target)
    return _translator_cache[key]

def get_file_type(filename):
    """Determine file type and processing method"""
    ext = filename.rsplit('.', 1)[1].lower()
    if ext == 'pdf':
        return 'pdf'
    elif ext in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif']:
        return 'image'
    elif ext in ['docx', 'doc']:
        return 'docx'
    elif ext in ['pptx', 'ppt']:
        return 'pptx'
    elif ext in ['xlsx', 'xls', 'csv']:
        return 'spreadsheet'
    elif ext == 'rtf':
        return 'rtf'
    elif ext in ['odt', 'ods', 'odp']:
        return 'opendocument'
    elif ext in ['html', 'htm']:
        return 'html'
    elif ext in ['xml', 'json']:
        return 'data'
    elif ext in ['txt', 'md']:
        return 'text'
    return None

def allowed_file(filename):
    """Check if file has allowed extension"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_page_count(pdf_path):
    """
    Get total page count using multiple methods (pdfinfo_from_path, then PyPDF2 fallback)
    """
    # Method 1: pdfinfo_from_path (fast, from poppler)
    try:
        pdf_info = pdfinfo_from_path(pdf_path)
        total_pages = pdf_info.get("Pages", pdf_info.get("pages"))
        logger.info(f"Page count from pdfinfo: {total_pages}")
        return total_pages
    except Exception as e:
        logger.warning(f"pdfinfo failed: {e}, trying PyPDF2")

    # Method 2: PyPDF2 (more reliable fallback)
    if PdfReader:
        try:
            reader = PdfReader(pdf_path)
            total_pages = len(reader.pages)
            logger.info(f"Page count from PyPDF2: {total_pages}")
            return total_pages
        except Exception as e:
            logger.error(f"PyPDF2 failed: {e}")

    logger.warning("Could not determine total page count")
    return None  # Unknown

def detect_pdf_language(pdf_path, sample_page=1):
    """
    Detect the primary language of a PDF by sampling the first page.
    Returns Tesseract language code (e.g., 'tam', 'eng').
    Uses OSD for non-Tamil scripts, verifies with character检查 for Tamil.
    """
    try:
        images = convert_from_path(
            pdf_path,
            dpi=150,
            first_page=sample_page,
            last_page=sample_page
        )
        if not images:
            logger.warning("No image from first page, defaulting to Tamil")
            return DEFAULT_LANG

        # Step 1: OSD (fast, good for non-Tamil scripts)
        osd_result = image_to_osd(images[0])
        logger.info(f"OSD result: {osd_result[:200]}...")

        script_match = re.search(r"Script: ([^\n]+)", osd_result)
        if script_match:
            script_name = script_match.group(1).strip()
            # If OSD says Tamil (or any non-Latin script), trust it
            if script_name != "Latin":
                lang_code = SCRIPT_TO_LANG.get(script_name, DEFAULT_LANG)
                logger.info(f"OSD detected: {script_name} → Language: {lang_code}")
                return lang_code

        # Step 2: OSD said Latin or failed — verify with quick Tamil OCR
        # Many Tamil PDFs are misdetected as Latin by OSD
        try:
            tam_text = pytesseract.image_to_string(images[0], lang='tam')
            tamil_chars = len(re.findall(r'[\u0B80-\u0BFF]', tam_text))
            logger.info(f"Tamil verification: {tamil_chars} Tamil chars found")
            if tamil_chars > 5:
                logger.info("Confirmed Tamil via character detection")
                return 'tam'
        except Exception:
            pass

        if script_match:
            script_name = script_match.group(1).strip()
            lang_code = SCRIPT_TO_LANG.get(script_name, DEFAULT_LANG)
            logger.info(f"No Tamil chars found, using OSD: {script_name} → {lang_code}")
            return lang_code

    except Exception as e:
        logger.warning(f"Language detection failed: {e}, defaulting to Tamil")
    return DEFAULT_LANG


def _ocr_page(img, lang, timeout, task_id, page):
    """Run OCR on a single image with timeout. Uses a fresh thread pool per call."""
    pool = ThreadPoolExecutor(max_workers=1)
    try:
        f = pool.submit(pytesseract.image_to_string, img, lang)
        result = f.result(timeout=timeout)
        return result
    except Exception as e:
        logger.warning(f"Task {task_id}: OCR failed on page {page} ({timeout}s timeout): {e}")
        raise
    finally:
        pool.shutdown(wait=False)


def process_pdf_ocr(pdf_path, lang=DEFAULT_LANG, dpi=200, task_id=None, output_file=None, start_page=None, end_page=None, translated_output_file=None, translate_to=None):
    """
    Process PDF pages in batches using OCR, write directly to file for memory efficiency.
    Updates progress_tracker if task_id is provided.
    Processes pages from start_page to end_page (inclusive) in batches of BATCH_SIZE.
    If end_page is None, processes only start_page (single page mode).
    If translated_output_file and translate_to are set, translates each page after OCR.
    Returns number of pages processed.
    """
    page_num = start_page if start_page else 1
    actual_end = end_page if end_page else page_num
    pages_processed = 0
    batch_size = BATCH_SIZE
    from_code = TESSERACT_TO_ISO.get(lang.split('+')[0] if '+' in lang else lang, 'en')

    current = page_num
    while current <= actual_end:
        # Pre-batch GC and memory check
        gc.collect()
        mem = _get_memory_mb()
        if mem > 400:
            logger.warning(f"Task {task_id}: RSS {mem}MB > 400MB before page {current}, forcing aggressive GC")
            gc.collect()
            gc.collect()
        
        # Check cancel
        if task_id:
            with progress_lock:
                if progress_tracker[task_id].get("cancelled"):
                    logger.info(f"Task {task_id}: OCR cancelled")
                    return pages_processed
        batch_end = min(current + batch_size - 1, actual_end)

        if task_id:
            with progress_lock:
                progress_tracker[task_id]["current_page"] = current

        convert_pool = ThreadPoolExecutor(max_workers=1)
        try:
            convert_future = convert_pool.submit(
                convert_from_path, pdf_path, dpi=dpi,
                first_page=current, last_page=batch_end
            )
            images = convert_future.result(timeout=CONVERT_TIMEOUT)
        except _CTimeoutError:
            logger.warning(f"Task {task_id}: Batch convert timed out pages {current}-{batch_end}, trying single-page fallback")
            images = None
        except Exception as e:
            logger.warning(f"Task {task_id}: Batch convert failed pages {current}-{batch_end}: {e}, trying single-page fallback")
            images = None
        finally:
            convert_pool.shutdown(wait=False)

        # If batch convert failed, process pages one at a time
        if images is None:
            single_page = current
            while single_page <= batch_end:
                # Check cancel
                if task_id:
                    with progress_lock:
                        if progress_tracker[task_id].get("cancelled"):
                            logger.info(f"Task {task_id}: OCR cancelled during fallback")
                            return pages_processed
                convert_pool2 = ThreadPoolExecutor(max_workers=1)
                try:
                    sf = convert_pool2.submit(
                        convert_from_path, pdf_path, dpi=dpi,
                        first_page=single_page, last_page=single_page
                    )
                    single_images = sf.result(timeout=CONVERT_TIMEOUT)
                except Exception:
                    logger.warning(f"Task {task_id}: Skipping page {single_page} (convert failed/timed out)")
                    single_page += 1
                    continue
                finally:
                    convert_pool2.shutdown(wait=False)

                if single_images:
                    page = single_page
                    t0 = time.time()
                    try:
                        text = _ocr_page(single_images[0], lang, OCR_TIMEOUT, task_id, page)
                        elapsed = time.time() - t0
                        if task_id:
                            with progress_lock:
                                times = progress_tracker[task_id].setdefault("page_times", [])
                                times.append(elapsed)
                                if len(times) > 5:
                                    times.pop(0)
                    except Exception:
                        text = f"[OCR failed on page {page}]"
                    single_images[0].close()
                    if output_file:
                        output_file.write(f"--- Page {page} ---\n{text}\n\n")
                        output_file.flush()
                    if translated_output_file and translate_to and not text.startswith("[OCR failed"):
                        translated_text = _translate_text(text, translate_to, source_lang=from_code)
                        if translated_text:
                            translated_output_file.write(f"{translated_text}\n\n")
                            translated_output_file.flush()
                    pages_processed += 1
                    if task_id:
                        with progress_lock:
                            progress_tracker[task_id]["current_page"] = page
                            total = progress_tracker[task_id].get("total_pages")
                            start_offset = progress_tracker[task_id].get("processing_start_page", 1)
                            if total:
                                relative_page = max(0, page - start_offset + 1)
                                pct = max(1, min(int((relative_page / total) * 100), 99))
                                progress_tracker[task_id]["percentage"] = pct
                    # Explicitly clean up single_images list
                    for sim in single_images:
                        try:
                            sim.close()
                        except Exception:
                            pass
                    single_images.clear()
                gc.collect()
                single_page += 1
            current = batch_end + 1
            gc.collect()
            continue

        if not images:
            break

        # OCR each page with timeout
        for i, img in enumerate(images):
            # Check cancel
            if task_id:
                with progress_lock:
                    if progress_tracker[task_id].get("cancelled"):
                        logger.info(f"Task {task_id}: OCR cancelled mid-batch")
                        img.close()
                        return pages_processed
            page = current + i
            t0 = time.time()
            try:
                text = _ocr_page(img, lang, OCR_TIMEOUT, task_id, page)
                elapsed = time.time() - t0
                if task_id:
                    with progress_lock:
                        times = progress_tracker[task_id].setdefault("page_times", [])
                        times.append(elapsed)
                        if len(times) > 5:
                            times.pop(0)
            except Exception:
                text = f"[OCR failed on page {page}]"
            img.close()
            if output_file:
                output_file.write(f"--- Page {page} ---\n{text}\n\n")
                output_file.flush()
            if translated_output_file and translate_to:
                translated_text = _translate_text(text, translate_to, source_lang=from_code)
                if translated_text:
                    translated_output_file.write(f"{translated_text}\n\n")
                    translated_output_file.flush()
            pages_processed += 1
            if task_id:
                    with progress_lock:
                        progress_tracker[task_id]["current_page"] = page
                        total = progress_tracker[task_id].get("total_pages")
                        start_offset = progress_tracker[task_id].get("processing_start_page", 1)
                        if total:
                            relative_page = max(0, page - start_offset + 1)
                            pct = max(1, min(int((relative_page / total) * 100), 99))
                            progress_tracker[task_id]["percentage"] = pct

        # Aggressive cleanup: close all remaining image handles
        for im in images:
            try:
                im.close()
            except Exception:
                pass
        del images
        gc.collect()
        # Log memory every 20 pages
        if pages_processed > 0 and pages_processed % 20 == 0:
            mem = _get_memory_mb()
            if mem:
                logger.info(f"Task {task_id}: Processed {pages_processed} pages, RSS ~{mem}MB")
        current = batch_end + 1

    return pages_processed

def process_image_file(image_path, lang='eng'):
    """Process image file with OCR"""
    try:
        img = Image.open(image_path)
        # Convert to RGB if necessary
        if img.mode != 'RGB':
            img = img.convert('RGB')
        text = pytesseract.image_to_string(img, lang=lang)
        return text
    except Exception as e:
        logger.error(f"Image processing error: {e}")
        raise

def process_docx_file(docx_path):
    """Extract text from DOCX file (no OCR needed)"""
    try:
        from docx import Document
        doc = Document(docx_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"DOCX processing error: {e}")
        raise

def process_pptx_file(pptx_path):
    """Extract text from PPTX file (no OCR needed)"""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        full_text = []
        for i, slide in enumerate(prs.slides, 1):
            full_text.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text.strip())
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"PPTX processing error: {e}")
        raise

def process_txt_file(txt_path):
    """Read text file directly"""
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encoding
        with open(txt_path, 'r', encoding='latin-1') as f:
            return f.read()

def process_spreadsheet_file(file_path, ext):
    """Extract text from Excel/CSV files"""
    try:
        import pandas as pd
        full_text = []
        
        if ext in ['xlsx', 'xls']:
            # Read all sheets
            xl_file = pd.ExcelFile(file_path)
            for sheet_name in xl_file.sheet_names:
                full_text.append(f"--- Sheet: {sheet_name} ---")
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                full_text.append(df.to_string(index=False))
        elif ext == 'csv':
            df = pd.read_csv(file_path)
            full_text.append(df.to_string(index=False))
        
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"Spreadsheet processing error: {e}")
        raise

def process_rtf_file(rtf_path):
    """Extract text from RTF files"""
    try:
        from pyth import open as pyth_open
        doc = pyth_open(rtf_path)
        return doc.get_text()
    except Exception as e:
        logger.error(f"RTF processing error: {e}")
        raise

def process_opendocument_file(odt_path, ext):
    """Extract text from OpenDocument files"""
    try:
        from odfpy import opendocument
        from xml.etree import ElementTree as ET
        
        doc = opendocument.load(odt_path)
        
        # Extract text from paragraphs
        full_text = []
        for para in doc.getElementsByType(odfpy.text.P):
            text = para.getFirstChildText()
            if text and text.strip():
                full_text.append(text.strip())
        
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"OpenDocument processing error: {e}")
        raise

def process_html_file(html_path):
    """Extract text from HTML files"""
    try:
        from bs4 import BeautifulSoup
        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        logger.error(f"HTML processing error: {e}")
        raise

def process_data_file(data_path, ext):
    """Extract text from XML/JSON files"""
    try:
        with open(data_path, 'r', encoding='utf-8') as f:
            if ext == 'json':
                import json
                data = json.load(f)
                return json.dumps(data, indent=2)
            else:  # XML
                from xml.etree import ElementTree as ET
                tree = ET.parse(f)
                root = tree.getroot()
                return ET.tostring(root, encoding='unicode', method='text')
    except Exception as e:
        logger.error(f"Data file processing error: {e}")
        raise

def process_file_background(task_id, file_path, filename, temp_dir, selected_lang='auto', start_page=1, end_page=None, dpi=200, file_type='pdf'):
    """
    Background thread function to process any file type and update progress
    """
    try:
        # Initialize progress
        with progress_lock:
            progress_tracker[task_id]["status"] = "processing"
            progress_tracker[task_id]["total_pages"] = 1  # Single unit for non-PDF
        
        # Create output file path
        output_filename = f"{os.path.splitext(filename)[0]}_ocr.txt"
        output_path = os.path.join(temp_dir, output_filename)
        
        with progress_lock:
            progress_tracker[task_id]["output_path"] = output_path
            progress_tracker[task_id]["output_filename"] = output_filename
            translate_to = progress_tracker[task_id].get("translate_to")
        
        # Create translated output file path if translate_to is set
        translated_output_path = None
        if translate_to:
            translated_output_path = os.path.join(temp_dir, f"{os.path.splitext(filename)[0]}_{translate_to}.txt")
            with progress_lock:
                progress_tracker[task_id]["translated_output_path"] = translated_output_path
        
        # Process based on file type
        if file_type == 'pdf':
            # Existing PDF processing logic
            # Determine language
            if selected_lang != 'auto':
                # Manual selection
                detected_lang = selected_lang
                logger.info(f"Task {task_id}: Using manually selected language {detected_lang}")
                with progress_lock:
                    progress_tracker[task_id]["status"] = "getting_page_count"
                    progress_tracker[task_id]["detected_language"] = detected_lang
            else:
                # Auto-detect
                with progress_lock:
                    progress_tracker[task_id]["status"] = "detecting_language"
                
                detected_lang = detect_pdf_language(file_path)
                logger.info(f"Task {task_id}: Detected language {detected_lang}")
                
                with progress_lock:
                    progress_tracker[task_id]["status"] = "getting_page_count"
                    progress_tracker[task_id]["detected_language"] = detected_lang
            
            logger.info(f"Task {task_id}: Starting page count for {file_path}")
            
            # Get total page count
            total_pages = get_page_count(file_path)
            logger.info(f"Task {task_id}: Total pages = {total_pages}")
            
            # Adjust page range
            actual_start = start_page
            actual_end = end_page
            
            if actual_end is None or actual_end > total_pages:
                actual_end = total_pages
            
            if actual_start > total_pages:
                actual_start = 1
            
            logger.info(f"Task {task_id}: Processing pages {actual_start} to {actual_end}")
            
            with progress_lock:
                progress_tracker[task_id]["status"] = "processing"
                page_range = actual_end - actual_start + 1 if actual_end else total_pages - actual_start + 1
                progress_tracker[task_id]["total_pages"] = page_range
                progress_tracker[task_id]["pdf_total_pages"] = total_pages
                progress_tracker[task_id]["processing_start_page"] = actual_start
            
            # Process PDF and write directly to file (memory efficient)
            logger.info(f"Task {task_id}: Starting OCR processing with language {detected_lang}")
            
            mega_ckpt = None
            originals_uploaded = False
            
            # Track start time
            ocr_start_time = time.time()
            pages_processed = 0
            last_checkpoint_pages = 0
            
            with open(output_path, 'w', encoding='utf-8') as output_file:
                t_out = None
                if translated_output_path:
                    t_out = open(translated_output_path, 'w', encoding='utf-8')
                current = actual_start
                try:
                    while actual_end is None or current <= actual_end:
                        if total_pages and current > total_pages:
                            break
                        
                        # Check for cancel
                        cancelled_flag = False
                        if task_id:
                            with progress_lock:
                                if progress_tracker[task_id].get("cancelled"):
                                    logger.info(f"Task {task_id}: Cancelled by user")
                                    progress_tracker[task_id]["status"] = "cancelled"
                                    progress_tracker[task_id]["error"] = "Cancelled by user"
                                    cancelled_flag = True
                        if cancelled_flag:
                            _save_progress(True)
                            return
                        
                        batch_end = min(current + BATCH_SIZE - 1, actual_end) if actual_end else current + BATCH_SIZE - 1
                        
                        # Process batch of pages
                        result = process_pdf_ocr(
                            file_path,
                            lang=detected_lang,
                            dpi=dpi,
                            task_id=task_id,
                            output_file=output_file,
                            start_page=current,
                            end_page=batch_end,
                            translated_output_file=t_out,
                            translate_to=translate_to
                        )
                        if result == 0:
                            break
                        
                        pages_processed += result
                        current = batch_end + 1
                        
                        # Log page progress every 10 pages
                        if pages_processed % 10 == 0:
                            logger.info(f"Task {task_id}: Processed {pages_processed} pages so far")
                            _ocr_keepalive_ping()
                        
                        # Save progress every page for crash recovery (throttled to max 1 write/2s)
                        _save_progress()
                        
                        # Upload checkpoint to Mega every CHECKPOINT_INTERVAL pages
                        if pages_processed - last_checkpoint_pages >= CHECKPOINT_INTERVAL:
                            last_checkpoint_pages = pages_processed
                            # Lazy Mega login (only when first checkpoint is due)
                            if mega_ckpt is None and os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD"):
                                mega_ckpt = init_mega()

                            if mega_ckpt:
                                # Upload original file on first checkpoint only (best-effort)
                                if not originals_uploaded:
                                    try:
                                        originals_handle = ensure_mega_folder(mega_ckpt, "ocr-originals")
                                        if originals_handle:
                                            res = mega_call(mega_ckpt, "upload", file_path, dest=originals_handle, dest_filename=filename, timeout=120)
                                            with progress_lock:
                                                progress_tracker[task_id]["mega_original_handle"] = str(getattr(res, 'node_id', res))
                                            originals_uploaded = True
                                            logger.info(f"Task {task_id}: Original uploaded to Mega for resume")
                                    except Exception as e:
                                        logger.warning(f"Task {task_id}: Original upload failed (checkpoint still saved): {e}")

                                checkpoint_data = {
                                    "task_id": task_id,
                                    "last_page": current - 1,
                                    "total_pages": total_pages,
                                    "filename": filename,
                                    "output_filename": output_filename,
                                    "detected_lang": detected_lang,
                                    "file_type": "pdf",
                                    "start_page": actual_start,
                                    "end_page": actual_end,
                                    "created_at": time.time(),
                                    "original_filename": filename
                                }
                                try:
                                    upload_checkpoint(mega_ckpt, task_id, output_path, checkpoint_data)
                                    logger.info(f"Task {task_id}: Checkpoint saved at page {current - 1}")
                                    with progress_lock:
                                        progress_tracker[task_id]["last_checkpoint_page"] = current - 1
                                    _save_progress()
                                except Exception as e:
                                    logger.warning(f"Task {task_id}: Checkpoint upload failed: {e}")
                finally:
                    if t_out:
                        t_out.close()
            
            logger.info(f"Task {task_id}: OCR returned {pages_processed} pages")
            
            # Check if cancelled before treating empty output as error
            if pages_processed == 0:
                cancelled_here = False
                with progress_lock:
                    if progress_tracker[task_id].get("cancelled"):
                        logger.info(f"Task {task_id}: Cancelled by user")
                        progress_tracker[task_id]["status"] = "cancelled"
                        progress_tracker[task_id]["error"] = "Cancelled by user"
                        cancelled_here = True
                if cancelled_here:
                    _save_progress(True)
                    return

            # Verify output file has content
            if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                logger.error(f"Task {task_id}: Output file empty or missing")
                with progress_lock:
                    progress_tracker[task_id]["status"] = "error"
                    progress_tracker[task_id]["error"] = "No text could be extracted"
                _save_progress(True)
                return
            
            # Update tracker with success (check cancel race)
            with progress_lock:
                if progress_tracker[task_id].get("cancelled"):
                    progress_tracker[task_id]["status"] = "cancelled"
                    progress_tracker[task_id]["error"] = "Cancelled by user"
                else:
                    progress_tracker[task_id]["status"] = "completed"
                    progress_tracker[task_id]["pages_processed"] = pages_processed
                    progress_tracker[task_id]["percentage"] = 100
            _save_progress(True)
            logger.info(f"Task {task_id}: OCR completed successfully")
            
        elif file_type == 'image':
            # Update language detection for images
            if selected_lang == 'auto':
                with progress_lock:
                    progress_tracker[task_id]["detected_language"] = 'eng'  # Default for images
            
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            # Process image with OCR
            text = process_image_file(file_path, lang=selected_lang if selected_lang != 'auto' else 'eng')
            
            # Write to output
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'docx':
            # Extract text directly (no OCR)
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            text = process_docx_file(file_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'pptx':
            # Extract text from slides
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            text = process_pptx_file(file_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'txt':
            # Read directly
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            text = process_txt_file(file_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'spreadsheet':
            # Extract text from spreadsheets
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            ext = filename.rsplit('.', 1)[1].lower()
            text = process_spreadsheet_file(file_path, ext)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'rtf':
            # Extract text from RTF
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            text = process_rtf_file(file_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'opendocument':
            # Extract text from OpenDocument
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            ext = filename.rsplit('.', 1)[1].lower()
            text = process_opendocument_file(file_path, ext)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'html':
            # Extract text from HTML
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            text = process_html_file(file_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'data':
            # Extract text from XML/JSON
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            ext = filename.rsplit('.', 1)[1].lower()
            text = process_data_file(file_path, ext)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
        
        # Save progress for non-PDF files before Mega upload attempt
        _save_progress()
        
        # Upload to Mega.nz cloud for permanent storage
        try:
            mega_link = upload_to_mega(output_path, output_filename)
            with progress_lock:
                progress_tracker[task_id]["download_link"] = mega_link
                if mega_link:
                    progress_tracker[task_id]["mega_uploaded"] = True
                    progress_tracker[task_id]["mega_status"] = "uploaded"
                    logger.info(f"Task {task_id}: Mega upload success - {mega_link}")
                    # Clean up checkpoint files on successful upload
                    try:
                        m_clean = init_mega()
                        if m_clean:
                            cleanup_checkpoints(m_clean, task_id)
                    except Exception:
                        pass
                else:
                    progress_tracker[task_id]["mega_uploaded"] = False
                    progress_tracker[task_id]["mega_status"] = "failed: upload_to_mega returned None"
                    logger.warning(f"Task {task_id}: Mega upload returned None (check Render logs)")
        except Exception as mega_err:
            error_msg = str(mega_err)
            logger.error(f"Task {task_id}: Mega upload error: {error_msg}")
            with progress_lock:
                progress_tracker[task_id]["mega_uploaded"] = False
                progress_tracker[task_id]["mega_status"] = f"failed: {error_msg}"

        # Persist output file to ocr-outputs folder
        persist_output(task_id)

        # Auto-translate to target language if set
        translated_link = _auto_translate_and_upload(task_id)
        if translated_link:
            with progress_lock:
                progress_tracker[task_id]["translated_link"] = translated_link

        # Mark as completed (respect cancel flag)
        with progress_lock:
            if not progress_tracker[task_id].get("cancelled"):
                progress_tracker[task_id]["status"] = "completed"
            progress_tracker[task_id]["completed_at"] = time.time()
            progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0)
        _save_progress(True)
            
    except Exception as e:
        logger.error(f"Task {task_id}: Error - {str(e)}")
        import traceback
        logger.error(f"Traceback: {traceback.format_exc()}")
        with progress_lock:
            progress_tracker[task_id]["status"] = "error"
            progress_tracker[task_id]["error"] = str(e)
        _save_progress(True)
    finally:
        _release_keepalive()

def cleanup_old_tasks():
    """Clean up completed tasks older than 24h. Preserves interrupted/error tasks so users can retry."""
    current_time = time.time()
    to_delete = []
    with progress_lock:
        for task_id, task in progress_tracker.items():
            status = task.get("status", "")
            # Never auto-delete interrupted, error, cancelled, or processing tasks
            if status in ("interrupted", "error", "cancelled", "cancelling", "processing", "resuming", "starting", "detecting_language", "getting_page_count"):
                continue
            if "created_at" not in task:
                continue
            # Only delete completed tasks older than 24 hours (increased from 1 hour)
            if current_time - task["created_at"] > 86400:  # 24 hours
                to_delete.append(task_id)

        for task_id in to_delete:
            logger.info(f"Cleaning up old completed task {task_id}")
            if progress_tracker[task_id].get("temp_dir"):
                shutil.rmtree(progress_tracker[task_id]["temp_dir"], ignore_errors=True)
            del progress_tracker[task_id]
    if to_delete:
        _save_progress()


def upload_to_mega(local_file_path, remote_filename):
    """Upload a file to Mega.nz ocr-outputs folder and return download link"""
    email = os.environ.get("MEGA_EMAIL")
    password = os.environ.get("MEGA_PWD")

    if not email or not password:
        logger.warning("MEGA_EMAIL/MEGA_PWD not set - skipping cloud upload")
        return None

    if not os.path.exists(local_file_path):
        logger.error(f"File not found for upload: {local_file_path}")
        return None

    from mega import Mega

    mega = Mega()
    m = mega.login(email, password)

    try:
        folders = mega_call(m, "find", "ocr-outputs")
    except Exception:
        logger.warning("Mega find failed during upload")
        return None
    if not folders:
        logger.info("Creating ocr-outputs folder in Mega")
        try:
            folder_node = mega_call(m, "create_folder", "ocr-outputs")
        except Exception:
            logger.warning("Mega create_folder failed during upload")
            return None
        dest = folder_node.get("ocr-outputs")
        if not dest:
            logger.error(f"Failed to get handle from create_folder result: {folder_node}")
            return None
    else:
        dest = folders[0] if isinstance(folders, (list, tuple)) else folders

    logger.info(f"Uploading {remote_filename} to Mega.nz...")
    logger.info(f"File size: {os.path.getsize(local_file_path)} bytes")
    try:
        file_node = mega_call(m, "upload", local_file_path, dest=dest, dest_filename=remote_filename, timeout=120)
    except Exception:
        logger.warning("Mega upload failed")
        return None
    logger.info(f"Upload completed, getting link...")
    try:
        link = mega_call(m, "get_upload_link", file_node)
    except Exception:
        logger.warning("Mega get_upload_link failed")
        return None
    logger.info(f"Mega.nz upload complete: {link}")
    return link


def rebuild_completed_from_mega():
    """Scan Mega ocr-outputs folder and restore completed tasks from _ocr.txt files"""
    email = os.environ.get("MEGA_EMAIL")
    password = os.environ.get("MEGA_PWD")
    if not email or not password:
        logger.info("MEGA_EMAIL/MEGA_PWD not set — skipping Mega restore scan")
        return

    try:
        from mega import Mega
        m = Mega().login(email, password)

        folder = mega_call(m, "find", "ocr-outputs", timeout=15)
        if isinstance(folder, (list, tuple)):
            folder = folder[0] if folder else None
        if not folder:
            logger.info("No ocr-outputs folder found in Mega — nothing to restore")
            return

        try:
            files = mega_call(m, "get_files_in_node", folder, timeout=15)
        except Exception as e:
            logger.warning(f"Mega get_files_in_node failed: {e}")
            return

        if not files:
            logger.info("Mega ocr-outputs folder is empty — nothing to restore")
            return

        now = time.time()
        restored = 0
        for nid, finfo in files.items():
            if not isinstance(finfo, dict):
                continue
            name = finfo.get('a', {}).get('n', '')
            if not name.endswith('_ocr.txt') or name.startswith('_'):
                continue

            hashlib = __import__('hashlib')
            tid = "mega_" + hashlib.md5(name.encode()).hexdigest()[:12]

            with progress_lock:
                if tid in progress_tracker:
                    continue

            orig_name = name[:-8].rstrip('_')  # Remove '_ocr.txt' + trailing underscore
            with progress_lock:
                progress_tracker[tid] = {
                    "current_page": 0, "status": "completed",
                    "result_path": None, "error": None,
                    "filename": orig_name, "output_filename": name,
                    "download_link": None,
                    "mega_node_id": nid,
                    "mega_uploaded": True, "mega_status": "uploaded",
                    "file_type": "pdf", "detected_language": "",
                    "pages_processed": 0, "percentage": 100,
                    "download_count": 0, "completed_at": now, "created_at": now
                }
                restored += 1

        logger.info(f"Mega restore scan complete: {restored} tasks restored")
    except Exception as e:
        logger.error(f"Failed to scan Mega for completed tasks: {e}")


def mega_call(m, method_name, *args, timeout=MEGA_LOGIN_TIMEOUT, **kwargs):
    """Call a Mega client method with timeout protection. Returns the method's result or raises."""
    fn = getattr(m, method_name)
    pool = ThreadPoolExecutor(max_workers=1)
    try:
        future = pool.submit(fn, *args, **kwargs)
        return future.result(timeout=timeout)
    except _CTimeoutError:
        logger.warning(f"Mega.{method_name} timed out after {timeout}s")
        raise
    finally:
        pool.shutdown(wait=False)


def init_mega():
    """Initialize and login to Mega.nz with timeout. Returns client or None."""
    email = os.environ.get("MEGA_EMAIL")
    password = os.environ.get("MEGA_PWD")
    if not email or not password:
        return None
    from concurrent.futures import ThreadPoolExecutor, TimeoutError as _TimeoutError
    try:
        from mega import Mega
        mega = Mega()
        ex = ThreadPoolExecutor(max_workers=1)
        try:
            future = ex.submit(mega.login, email, password)
            return future.result(timeout=MEGA_LOGIN_TIMEOUT)
        finally:
            ex.shutdown(wait=False)
    except _TimeoutError:
        logger.warning("Mega login timed out (network issue?)")
        return None
    except Exception as e:
        logger.error(f"Mega login failed: {e}")
        return None


def ensure_mega_folder(m, folder_name):
    """Find Mega folder, create if needed. Return folder handle string."""
    try:
        folder = mega_call(m, "find", folder_name)
    except Exception:
        logger.warning(f"Mega find failed for folder {folder_name}")
        return None
    if folder:
        return folder[0]
    try:
        result = mega_call(m, "create_folder", folder_name)
    except Exception:
        logger.warning(f"Mega create_folder failed for {folder_name}")
        return None
    return result.get(folder_name)


def upload_checkpoint(m, task_id, output_path, metadata):
    """Upload checkpoint (partial output + metadata) to ocr-checkpoints/ folder in Mega."""
    import json, tempfile as _tf
    folder_name = "ocr-checkpoints"
    folder_handle = ensure_mega_folder(m, folder_name)
    if not folder_handle:
        logger.error(f"Cannot access/create {folder_name} in Mega")
        return False

    # Delete old checkpoint file if exists
    try:
        old_ckpt = mega_call(m, "find", f"{folder_name}/{task_id}.checkpoint")
        if old_ckpt:
            mega_call(m, "delete", old_ckpt[0])
    except Exception:
        pass
    # Delete old output file if exists
    try:
        old_out = mega_call(m, "find", f"{folder_name}/{task_id}_output.txt")
        if old_out:
            mega_call(m, "delete", old_out[0])
    except Exception:
        pass

    # Upload metadata JSON first
    ckpt_file = _tf.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8')
    json.dump(metadata, ckpt_file)
    ckpt_file.close()
    try:
        mega_call(m, "upload", ckpt_file.name, dest=folder_handle, dest_filename=f"{task_id}.checkpoint", timeout=120)
    except Exception as e:
        logger.error(f"Checkpoint metadata upload failed: {e}")
        os.unlink(ckpt_file.name)
        return False
    os.unlink(ckpt_file.name)

    # Then upload partial output file
    try:
        mega_call(m, "upload", output_path, dest=folder_handle, dest_filename=f"{task_id}_output.txt", timeout=120)
    except Exception as e:
        logger.error(f"Checkpoint output upload failed: {e}")
        return False

    logger.info(f"Checkpoint saved for task {task_id} at page {metadata.get('last_page')}")
    return True


def cleanup_checkpoints(m, task_id):
    """Delete checkpoint files for a completed task."""
    folder_name = "ocr-checkpoints"
    try:
        old_ckpt = mega_call(m, "find", f"{folder_name}/{task_id}.checkpoint")
        if old_ckpt:
            mega_call(m, "delete", old_ckpt[0])
    except Exception:
        pass
    try:
        old_out = mega_call(m, "find", f"{folder_name}/{task_id}_output.txt")
        if old_out:
            mega_call(m, "delete", old_out[0])
    except Exception:
        pass
    logger.info(f"Checkpoint files cleaned up for task {task_id}")


def scan_and_resume_checkpoints():
    """Called at app startup. Scans Mega for checkpoint files and resumes tasks."""
    logger.info("Scanning for incomplete tasks to resume...")
    try:
        m = init_mega()
        if not m:
            logger.info("Mega not configured - skipping resume scan")
            return

        folder_handle = ensure_mega_folder(m, "ocr-checkpoints")
        if not folder_handle:
            logger.info("No ocr-checkpoints folder - nothing to resume")
            return

        files_in_folder = mega_call(m, "get_files_in_node", folder_handle)
        if not files_in_folder:
            logger.info("No checkpoint files found")
            return

        import json, tempfile as _tf

        for handle, node in files_in_folder.items():
            name = node.get('a', {}).get('n', '')
            if not name.endswith('.checkpoint'):
                continue
            task_id = name[:-len('.checkpoint')]

            logger.info(f"Found incomplete task {task_id}, attempting resume...")
            try:
                temp_dir = _tf.mkdtemp()
                mega_call(m, "download", (handle, node), dest_path=temp_dir)

                ckpt_path = os.path.join(temp_dir, name)
                with open(ckpt_path, 'r', encoding='utf-8') as f:
                    metadata = json.load(f)

                output_name = f"{task_id}_output.txt"
                output_node = mega_call(m, "find", f"ocr-checkpoints/{output_name}")
                if not output_node:
                    logger.warning(f"Task {task_id}: Partial output not found, skipping")
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    continue
                mega_call(m, "download", output_node, dest_path=temp_dir)

                original_name = metadata.get("original_filename", "")
                original_dl_path = None
                if original_name:
                    orig_node = mega_call(m, "find", f"ocr-originals/{original_name}")
                    if orig_node:
                        mega_call(m, "download", orig_node, dest_path=temp_dir)
                        original_dl_path = os.path.join(temp_dir, original_name)

                last_page = metadata.get("last_page", 0)
                total_pages = metadata.get("total_pages", 0)
                resume_from = last_page + 1

                if not original_dl_path:
                    logger.warning(f"Task {task_id}: Original file not found, cannot resume")
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    continue

                if resume_from > total_pages:
                    logger.info(f"Task {task_id}: Already complete (page {last_page}/{total_pages})")
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    continue

                output_filename = metadata.get("output_filename", f"{os.path.splitext(original_name)[0]}_ocr.txt")
                partial_path = os.path.join(temp_dir, output_name)
                new_output_path = os.path.join(temp_dir, output_filename)

                if os.path.exists(partial_path):
                    os.rename(partial_path, new_output_path)

                detected_lang = metadata.get("detected_lang", "tam")
                file_type = metadata.get("file_type", "pdf")
                actual_start = metadata.get("start_page", 1)
                actual_end = metadata.get("end_page", total_pages)

                with progress_lock:
                    if task_id in progress_tracker:
                        logger.info(f"Task {task_id} already in tracker, skipping")
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        continue
                    progress_tracker[task_id] = {
                        "current_page": resume_from,
                        "status": "resuming",
                        "result_path": None,
                        "error": None,
                        "filename": metadata.get("filename", original_name),
                        "temp_dir": temp_dir,
                        "total_pages": total_pages,
                        "percentage": int((last_page / total_pages) * 100) if total_pages else 0,
                        "detected_language": detected_lang,
                        "created_at": time.time(),
                        "file_type": file_type,
                        "output_path": new_output_path,
                        "output_filename": output_filename,
                        "resumed_from_page": resume_from
                    }

                logger.info(f"Resuming task {task_id} from page {resume_from}/{total_pages}")
                thread = threading.Thread(
                    target=resume_ocr_processing,
                    args=(task_id, original_dl_path, metadata, new_output_path, temp_dir)
                )
                thread.daemon = True
                thread.start()
                _ensure_keepalive()

            except Exception as e:
                logger.error(f"Failed to resume task {task_id}: {e}")
                import traceback
                logger.error(traceback.format_exc())

    except Exception as e:
        logger.error(f"Checkpoint scan error: {e}")


def resume_ocr_processing(task_id, original_path, metadata, output_path, temp_dir):
    """Resume OCR from last checkpoint page."""
    try:
        last_page = metadata.get("last_page", 0)
        total_pages = metadata.get("total_pages", 0)
        detected_lang = metadata.get("detected_lang", "tam")
        actual_start = metadata.get("start_page", 1)
        actual_end = metadata.get("end_page", total_pages)
        filename = metadata.get("filename", "document.pdf")
        output_filename = metadata.get("output_filename", "document_ocr.txt")

        resume_from = last_page + 1
        if resume_from > actual_end or (total_pages and resume_from > total_pages):
            with progress_lock:
                progress_tracker[task_id]["status"] = "completed"
                progress_tracker[task_id]["percentage"] = 100
            logger.info(f"Task {task_id}: Already complete, nothing to resume")
            return

        with progress_lock:
            progress_tracker[task_id]["status"] = "processing"
            progress_tracker[task_id]["total_pages"] = total_pages
        mega_ckpt = init_mega() if os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD") else None

        pages_processed = last_page
        last_checkpoint_pages = last_page

        with open(output_path, 'a', encoding='utf-8') as output_file:
            current = resume_from
            while actual_end is None or current <= actual_end:
                if total_pages and current > total_pages:
                    break

                batch_end = min(current + BATCH_SIZE - 1, actual_end) if actual_end else current + BATCH_SIZE - 1

                result = process_pdf_ocr(
                    original_path,
                    lang=detected_lang,
                    dpi=200,
                    task_id=task_id,
                    output_file=output_file,
                    start_page=current,
                    end_page=batch_end
                )
                if result == 0:
                    break

                pages_processed += result
                current = batch_end + 1

                if mega_ckpt and (pages_processed - last_checkpoint_pages >= CHECKPOINT_INTERVAL):
                    last_checkpoint_pages = pages_processed
                    try:
                        checkpoint_data = {
                            "task_id": task_id,
                            "last_page": current - 1,
                            "total_pages": total_pages,
                            "filename": filename,
                            "output_filename": output_filename,
                            "detected_lang": detected_lang,
                            "file_type": "pdf",
                            "start_page": actual_start,
                            "end_page": actual_end,
                            "created_at": time.time(),
                            "original_filename": filename
                        }
                        upload_checkpoint(mega_ckpt, task_id, output_path, checkpoint_data)
                    except Exception as e:
                        logger.warning(f"Task {task_id}: Resume checkpoint failed: {e}")

        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            logger.error(f"Task {task_id}: Resumed output empty")
            with progress_lock:
                progress_tracker[task_id]["status"] = "error"
                progress_tracker[task_id]["error"] = "Resumed OCR produced no output"
            return

        try:
            mega_link = upload_to_mega(output_path, output_filename)
            with progress_lock:
                progress_tracker[task_id]["download_link"] = mega_link
                if mega_link:
                    progress_tracker[task_id]["mega_uploaded"] = True
                    progress_tracker[task_id]["mega_status"] = "uploaded"
                    if mega_ckpt:
                        try:
                            cleanup_checkpoints(mega_ckpt, task_id)
                        except Exception:
                            pass
                else:
                    progress_tracker[task_id]["mega_uploaded"] = False
                    progress_tracker[task_id]["mega_status"] = "failed: upload_to_mega returned None"
        except Exception as mega_err:
            logger.error(f"Task {task_id}: Resume final upload error: {mega_err}")
            with progress_lock:
                progress_tracker[task_id]["mega_uploaded"] = False
                progress_tracker[task_id]["mega_status"] = f"failed: {mega_err}"

        persist_output(task_id)

        # Auto-translate to target language if set
        translated_link = _auto_translate_and_upload(task_id)
        if translated_link:
            with progress_lock:
                progress_tracker[task_id]["translated_link"] = translated_link

        with progress_lock:
            progress_tracker[task_id]["status"] = "completed"
            progress_tracker[task_id]["pages_processed"] = pages_processed
            progress_tracker[task_id]["percentage"] = 100
            progress_tracker[task_id]["completed_at"] = time.time()
            progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0)
        logger.info(f"Task {task_id}: Resume completed ({pages_processed} pages)")

    except Exception as e:
        logger.error(f"Task {task_id}: Resume error: {e}")
        import traceback
        logger.error(traceback.format_exc())
        with progress_lock:
            progress_tracker[task_id]["status"] = "error"
            progress_tracker[task_id]["error"] = str(e)
    finally:
        _release_keepalive()


def _ensure_keepalive():
    global _active_tasks, _keepalive_thread
    with _keepalive_lock:
        _active_tasks += 1
        if _keepalive_thread is None or not _keepalive_thread.is_alive():
            _keepalive_thread = threading.Thread(target=_keepalive_loop, daemon=True)
            _keepalive_thread.start()
            logger.info("Keepalive thread started")


def _release_keepalive():
    global _active_tasks
    with _keepalive_lock:
        _active_tasks -= 1
        if _active_tasks < 0:
            _active_tasks = 0


def _keepalive_loop():
    render_url = os.environ.get("RENDER_EXTERNAL_URL") or os.environ.get("RENDER_URL", "")
    local_url = f"http://localhost:{os.environ.get('PORT', 8080)}"

    while True:
        with _keepalive_lock:
            if _active_tasks <= 0:
                break
        success = False
        if render_url:
            try:
                requests.get(f"{render_url}/health", timeout=15)
                success = True
            except Exception:
                pass
        if not success:
            try:
                requests.get(f"{local_url}/health", timeout=5)
                success = True
            except Exception:
                pass
        if not success:
            logger.warning("Keepalive: all endpoints unreachable")
        time.sleep(60)  # Ping every 60s instead of 120s for faster response


def _ocr_keepalive_ping():
    """Secondary keepalive: called from within the OCR thread itself.
    Pings localhost every 60s so Render doesn't sleep even if the main
    keepalive thread is delayed or dead."""
    try:
        local_url = f"http://localhost:{os.environ.get('PORT', 8080)}"
        requests.get(f"{local_url}/health", timeout=5)
    except Exception:
        pass


@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        # Check if file was uploaded
        if 'file' not in request.files:
            flash('No file selected')
            return redirect(request.url)

        file = request.files['file']

        if file.filename == '':
            flash('No file selected')
            return redirect(request.url)

        if not allowed_file(file.filename):
            flash('Unsupported file type')
            return redirect(request.url)
        
        # Detect file type
        file_type = get_file_type(file.filename)
        if not file_type:
            flash('Unsupported file type')
            return redirect(request.url)
        
        # Get language selection (only used for PDF and image)
        selected_lang = request.form.get('language', 'auto')
        logger.info(f"Language selection: {selected_lang}")
        
        # Get target translation language (auto-translate after OCR)
        translate_to = request.form.get('translate_to', 'none')
        logger.info(f"Translate to: {translate_to}")
        
        # Get page range (only for PDF)
        page_range = request.form.get('page_range', 'all')
        start_page = 1
        end_page = None
        
        if page_range == 'custom':
            try:
                start_page = int(request.form.get('start_page', 1))
                end_page_str = request.form.get('end_page', '').strip()
                if end_page_str:
                    end_page = int(end_page_str)
                logger.info(f"Custom page range: {start_page} to {end_page}")
            except ValueError:
                logger.warning("Invalid page range, using defaults")
                start_page = 1
                end_page = None

        # Save uploaded file to temp location
        filename = secure_filename(file.filename)
        temp_dir = tempfile.mkdtemp()
        file_path = os.path.join(temp_dir, filename)
        file.save(file_path)
        
        # Verify file was saved
        if not os.path.exists(file_path):
            logger.error(f"Failed to save file to {file_path}")
            flash('Failed to save uploaded file')
            return redirect(request.url)
        
        logger.info(f"File saved to {file_path}, size: {os.path.getsize(file_path)} bytes")
        
        # Create a task ID for progress tracking
        task_id = str(uuid.uuid4())
        logger.info(f"Created task {task_id} for file {filename}")
        
        # Initialize progress tracker
        with progress_lock:
            progress_tracker[task_id] = {
                "current_page": 0,
                "status": "starting",
                "result_path": None,
                "error": None,
                "filename": filename,
                "temp_dir": temp_dir,
                "file_path": file_path,
                "total_pages": None,
                "percentage": 0,
                "detected_language": selected_lang if selected_lang != 'auto' else None,
                "selected_lang": selected_lang,
                "start_page": start_page,
                "end_page": end_page,
                "created_at": time.time(),
                "file_type": file_type,
                "cancelled": False,
                "translate_to": translate_to if translate_to != 'none' else None
            }
        _save_progress(True)
        
        # Start background processing thread
        logger.info(f"Starting background thread for task {task_id}, type: {file_type}")
        thread = threading.Thread(
            target=process_file_background,
            args=(task_id, file_path, filename, temp_dir, selected_lang, start_page, end_page, 200, file_type)
        )
        thread.daemon = True
        thread.start()
        logger.info(f"Background thread started for task {task_id}")
        _ensure_keepalive()
        
        # Render the processing page with task ID
        return render_template('processing.html', task_id=task_id)
    
    return render_template('index.html')

@app.route('/progress/<task_id>')
def get_progress(task_id):
    """Return progress status as JSON"""
    try:
        with progress_lock:
            if task_id not in progress_tracker:
                logger.warning(f"Progress check: task {task_id} not found in tracker")
                return jsonify({"status": "not_found"}), 404

            task = progress_tracker[task_id]
            
            # Get language code - use selected language if available
            lang_code = task.get("detected_language")
            if not lang_code or lang_code == 'auto':
                lang_code = DEFAULT_LANG
            
            lang_name = LANG_CODE_TO_NAME.get(lang_code, lang_code)
            
            # Compute ETA
            eta = None
            if task["status"] in ("processing", "resuming"):
                page_times = task.get("page_times")
                total = task.get("total_pages")
                current_page = task.get("current_page", 0)
                start_page_offset = task.get("processing_start_page", 1)
                if page_times and total and current_page > 0:
                    avg = sum(page_times) / len(page_times)
                    pages_done = max(0, current_page - start_page_offset + 1)
                    remaining = max(0, total - pages_done)
                    if remaining > 0:
                        eta = int(avg * remaining)
            
            return jsonify({
                "status": task["status"],
                "current_page": task["current_page"],
                "total_pages": task.get("total_pages"),
                "pdf_total_pages": task.get("pdf_total_pages"),
                "percentage": task.get("percentage", 0),
                "error": task["error"],
                "filename": task["filename"],
                "detected_language": lang_code,
                "language_name": lang_name,
                "file_type": task.get("file_type", "pdf"),
                "eta": eta
            })
    except Exception as e:
        logger.error(f"Progress endpoint error for {task_id}: {e}")
        return jsonify({"status": "error", "error": "Failed to get progress"}), 500


@app.route('/check/<task_id>')
def check_progress_page(task_id):
    """Styled HTML progress page for checking task status."""
    return render_template('progress_check.html', task_id=task_id)


@app.route('/download/<task_id>')
def download_result(task_id):
    """Download OCR result file"""
    with progress_lock:
        if task_id not in progress_tracker:
            flash('Task not found')
            return redirect(url_for('index'))

        task = progress_tracker[task_id]

        if task["status"] != "completed":
            flash('Processing not completed')
            return redirect(url_for('index'))

        if not task["output_path"] or not os.path.exists(task["output_path"]):
            flash('Result file not found')
            return redirect(url_for('index'))

        # Send file as download
        response = send_file(
            task["output_path"],
            as_attachment=True,
            download_name=task["output_filename"],
            mimetype='text/plain'
        )

        # Track download count
        with progress_lock:
            progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0) + 1

        # Cleanup after sending (optional, can keep for debugging)
        # shutil.rmtree(task["temp_dir"], ignore_errors=True)
        # with progress_lock:
        #     del progress_tracker[task_id]

        return response


@app.route('/download-docx/<task_id>')
def download_docx(task_id):
    """Download OCR result as a Word document."""
    with progress_lock:
        if task_id not in progress_tracker:
            flash('Task not found')
            return redirect(url_for('index'))
        task = progress_tracker[task_id]
        if task["status"] != "completed":
            flash('Processing not completed')
            return redirect(url_for('index'))
        if not task["output_path"] or not os.path.exists(task["output_path"]):
            flash('Result file not found')
            return redirect(url_for('index'))

    from docx import Document
    from io import BytesIO

    try:
        with open(task["output_path"], 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception:
        flash('Error reading result file')
        return redirect(url_for('index'))

    doc = Document()
    for para in content.split('\n\n'):
        stripped = para.strip()
        if stripped:
            doc.add_paragraph(stripped)

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)

    docx_filename = task.get("output_filename", "output.txt")
    docx_filename = docx_filename.rsplit('.', 1)[0] + '.docx'

    with progress_lock:
        progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0) + 1

    return send_file(
        buf,
        as_attachment=True,
        download_name=docx_filename,
        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    )


@app.route('/track-download/<task_id>', methods=['POST'])
def track_download(task_id):
    """Increment download count for cloud (Mega) link clicks"""
    with progress_lock:
        if task_id in progress_tracker:
            progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0) + 1
            return jsonify({"ok": True}), 200
    return jsonify({"ok": False}), 404


@app.route('/mega-link/<task_id>', methods=['POST'])
def get_mega_link(task_id):
    """Generate and cache a Mega download link for a task on demand."""
    with progress_lock:
        task = progress_tracker.get(task_id)
        if not task:
            return jsonify({"error": "Not found"}), 404
        if task.get("download_link"):
            return jsonify({"link": task["download_link"]}), 200
        nid = task.get("mega_node_id")
        if not nid:
            return jsonify({"error": "No Mega node"}), 404
    try:
        from mega import Mega
        email = os.environ.get("MEGA_EMAIL")
        password = os.environ.get("MEGA_PWD")
        if not email or not password:
            return jsonify({"error": "Mega not configured"}), 500
        m = Mega().login(email, password)
        link = mega_call(m, "get_link", nid, timeout=30)
        if not link:
            return jsonify({"error": "Failed to get link"}), 500
        with progress_lock:
            task = progress_tracker.get(task_id)
            if task:
                task["download_link"] = link
        return jsonify({"link": link}), 200
    except Exception as e:
        logger.error(f"Failed to get Mega link for {task_id}: {e}")
        return jsonify({"error": "Failed to get download link"}), 500


@app.route('/preview/<task_id>')
def preview_text(task_id):
    """Return first 1000 chars of OCR result for preview"""
    with progress_lock:
        task = progress_tracker.get(task_id)
        if not task or task.get("status") != "completed":
            return jsonify({"error": "Not available"}), 404
        output_path = task.get("output_path")
    if not output_path or not os.path.exists(output_path):
        return jsonify({"text": "File not found"})
    try:
        with open(output_path, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read(1000)
        return jsonify({"text": text})
    except Exception:
        return jsonify({"text": "Error reading file"})


@app.route('/translate', methods=['POST'])
def translate_text():
    """Translate full OCR result using Argos Translate (offline, free) and save to Mega ocr-translated folder."""
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "No data provided"}), 400
    target = data.get("target_lang", "en")
    source = data.get("source_lang", "en")
    task_id = data.get("task_id")

    if not task_id:
        return jsonify({"error": "task_id required"}), 400

    if not _DEEP_TRANSLATOR_AVAILABLE:
        return jsonify({"error": "Translation library (deep-translator) not installed"}), 500

    # Read the full output file
    with progress_lock:
        task = progress_tracker.get(task_id)
        if not task or task.get("status") != "completed":
            return jsonify({"error": "Task not completed or not found"}), 404
        output_path = task.get("output_path")
        base_filename = task.get("output_filename", "output.txt")
    if not output_path or not os.path.exists(output_path):
        return jsonify({"error": "Output file not found"}), 404

    try:
        with open(output_path, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
    except Exception:
        return jsonify({"error": "Failed to read output file"}), 500

    if not text.strip():
        return jsonify({"error": "Output file is empty"}), 400

    try:
        t = _get_translator(source, target)
        translated = t.translate(text)
    except Exception:
        return jsonify({"error": "Translation failed - text may be too large"}), 500

    if not translated:
        return jsonify({"error": "Translation returned empty"}), 500

    # Save translated text to Mega ocr-translated folder
    mega_link = None
    try:
        translated_filename = base_filename.rsplit('.', 1)[0] + f"_{target}.txt"
        import tempfile as _tf
        tmp = _tf.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8')
        tmp.write(translated)
        tmp.close()
        m = _get_mega_client()
        if m:
            folder_handle = ensure_mega_folder(m, "ocr-translated")
            if folder_handle:
                mega_call(m, "upload", tmp.name, dest=folder_handle, dest_filename=translated_filename, timeout=120)
                uploaded_node = mega_call(m, "find", f"ocr-translated/{translated_filename}")
                if uploaded_node:
                    mega_link = mega_call(m, "get_link", uploaded_node[0], timeout=30)
        os.unlink(tmp.name)
    except Exception as e:
        logger.error(f"Failed to upload translated file to Mega for {task_id}: {e}")

    resp_data = {"translated": translated}
    if mega_link:
        resp_data["mega_link"] = mega_link
    return jsonify(resp_data)


@app.route('/cancel/<task_id>', methods=['POST'])
def cancel_task(task_id):
    """Cancel a running OCR task"""
    with progress_lock:
        task = progress_tracker.get(task_id)
        if not task:
            return jsonify({"error": "Task not found"}), 404
        if task["status"] not in ("processing", "resuming", "starting", "detecting_language", "getting_page_count"):
            return jsonify({"error": "Task is not running"}), 400
        task["cancelled"] = True
        task["status"] = "cancelling"
    _save_progress(True)
    logger.info(f"Task {task_id}: Cancel requested")
    return jsonify({"status": "cancelling"}), 200


@app.route('/retry/<task_id>', methods=['POST'])
def retry_task(task_id):
    """Retry a failed OCR task"""
    with progress_lock:
        old = progress_tracker.get(task_id)
        if not old:
            return jsonify({"error": "Task not found"}), 404
        if old.get("status") not in ("error", "interrupted", "completed_with_error"):
            return jsonify({"error": "Only failed or interrupted tasks can be retried"}), 400

        filename = old.get("filename", "input.pdf")
        file_type = old.get("file_type", "pdf")
        temp_dir = old.get("temp_dir")
        file_path = old.get("file_path")
        selected_lang = old.get("selected_lang", "auto")
        last_ckpt = old.get("last_checkpoint_page", 0)
        current_page = old.get("current_page", 0)
        end_page = old.get("end_page")
        old_total = old.get("total_pages")
        mega_node = old.get("mega_original_handle")
        detected_lang = old.get("detected_language")
        translate_to = old.get("translate_to")

    # Everything after here is outside the lock

    # Try to locate the original file
    src = file_path or (os.path.join(temp_dir, filename) if temp_dir else None)
    needs_new_temp = False

    if not src or not os.path.exists(src):
        if mega_node:
            new_temp = tempfile.mkdtemp()
            src = os.path.join(new_temp, filename)
            try:
                m = init_mega()
                mega_call(m, "download", mega_node, dest_path=src, timeout=120)
                needs_new_temp = True
                logger.info(f"Task {task_id}: Downloaded original from Mega for retry")
            except Exception as e:
                shutil.rmtree(new_temp, ignore_errors=True)
                return jsonify({"error": "Failed to download original from cloud"}), 400
        else:
            return jsonify({"error": "Original file not found and no Mega backup"}), 400

    # Create new task
    new_id = str(uuid.uuid4())
    new_temp_dir = tempfile.mkdtemp() if needs_new_temp else temp_dir
    if needs_new_temp:
        shutil.copy2(src, os.path.join(new_temp_dir, filename))
        src = os.path.join(new_temp_dir, filename)
    else:
        src = file_path

    # Resume from last known page, not from page 1
    resume_start = (last_ckpt or current_page or 0) + 1
    if resume_start < 1:
        resume_start = 1

    with progress_lock:
        progress_tracker[new_id] = {
            "current_page": resume_start - 1,
            "status": "starting",
            "result_path": None,
            "error": None,
            "filename": filename,
            "temp_dir": new_temp_dir,
            "file_path": src,
            "total_pages": old_total,
            "percentage": 0,
            "detected_language": detected_lang,
            "selected_lang": selected_lang,
            "start_page": resume_start,
            "end_page": end_page,
            "created_at": time.time(),
            "file_type": file_type,
            "cancelled": False,
            "retry_of": task_id,
            "translate_to": translate_to
        }
    _save_progress(True)

    thread = threading.Thread(
        target=process_file_background,
        args=(new_id, src, filename, new_temp_dir, selected_lang, resume_start, end_page, 200, file_type)
    )
    thread.daemon = True
    thread.start()
    _ensure_keepalive()

    return jsonify({"task_id": new_id}), 200


@app.route('/health')
def health():
    """Health check endpoint for Render"""
    try:
        version = pytesseract.get_tesseract_version()
        # Check if Tamil and English packs are available
        languages = pytesseract.get_languages()
        return {
            'status': 'healthy',
            'tesseract_version': str(version),
            'languages_available': languages
        }, 200
    except Exception as e:
        return {'status': 'unhealthy', 'error': str(e)}, 500


@app.errorhandler(413)
def too_large(e):
    flash('File too large. Maximum size is 100MB.')
    return redirect(url_for('index'))


@app.route('/clear-downloads', methods=['POST'])
def clear_downloads():
    data = request.get_json(silent=True) or {}
    task_ids = data.get("task_ids")
    cleared = 0
    with progress_lock:
        if task_ids:
            for tid in task_ids:
                if tid in progress_tracker:
                    task = progress_tracker[tid]
                    op = task.get("output_path")
                    if op and os.path.exists(op):
                        try: os.remove(op)
                        except: pass
                    td = task.get("temp_dir")
                    if td:
                        shutil.rmtree(td, ignore_errors=True)
                    del progress_tracker[tid]
                    cleared += 1
        else:
            for tid in list(progress_tracker.keys()):
                task = progress_tracker[tid]
                op = task.get("output_path")
                if op and os.path.exists(op):
                    try: os.remove(op)
                    except: pass
                td = task.get("temp_dir")
                if td:
                    shutil.rmtree(td, ignore_errors=True)
                del progress_tracker[tid]
                cleared += 1
    if cleared:
        _save_progress(True)
    return jsonify({"cleared": cleared}), 200


@app.route('/downloads')
def downloads_page():
    all_tasks = []
    restoring = not _mega_restore_done.is_set()
    with progress_lock:
        for task_id, task in progress_tracker.items():
            status = task.get("status", "")

            # Compute ETA for in-progress tasks
            eta = None
            if status in ("processing", "resuming"):
                page_times = task.get("page_times")
                total = task.get("total_pages")
                current_page = task.get("current_page", 0)
                start_page_offset = task.get("processing_start_page", 1)
                if page_times and total and current_page > 0:
                    avg = sum(page_times) / len(page_times)
                    pages_done = max(0, current_page - start_page_offset + 1)
                    remaining = max(0, total - pages_done)
                    if remaining > 0:
                        eta = int(avg * remaining)

            info = {
                "task_id": task_id,
                "filename": task.get("output_filename", task.get("filename", "Unknown")),
                "download_link": task.get("download_link", ""),
                "language": task.get("detected_language", ""),
                "file_type": task.get("file_type", ""),
                "pages_processed": task.get("pages_processed", 0),
                "mega_uploaded": task.get("mega_uploaded", False),
                "mega_status": task.get("mega_status", ""),
                "completed_at": task.get("completed_at", 0),
                "created_at": task.get("created_at", 0),
                "download_count": task.get("download_count", 0),
                "translated_link": task.get("translated_link", ""),
                "translate_to": task.get("translate_to", ""),
                "status": status,
                "percentage": task.get("percentage", 0),
                "eta": eta,
                "last_checkpoint_page": task.get("last_checkpoint_page")
            }
            all_tasks.append(info)
    all_tasks.reverse()
    return render_template("downloads.html", downloads=all_tasks, restoring=restoring)



# Fast local restore (synchronous, <1s), Mega scan in background (fast with deferred links)
os.makedirs(OUTPUT_DIR, exist_ok=True)
saved = _load_progress()
if saved:
    for tid, data in saved.items():
        if data.get("status") not in ("completed", "error", "cancelled"):
            data["status"] = "interrupted"
        progress_tracker[tid] = data
    logger.info(f"Restored {len(saved)} persisted tasks from {PROGRESS_FILE}")

rebuild_completed_from_local()

with progress_lock:
    for tid, data in progress_tracker.items():
        if data.get("status") == "completed":
            op = data.get("output_path")
            if op and not os.path.exists(op):
                out_fn = data.get("output_filename")
                alt = os.path.join(OUTPUT_DIR, f"{tid}_{out_fn}") if out_fn else None
                if alt and os.path.exists(alt):
                    data["output_path"] = alt

if os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD"):
    def _mega_background():
        try:
            rebuild_completed_from_mega()
            scan_and_resume_checkpoints()
        except Exception:
            logger.error("Mega background restore failed", exc_info=True)
        finally:
            _mega_restore_done.set()
    threading.Thread(target=_mega_background, daemon=True).start()
else:
    _mega_restore_done.set()

# Periodic cleanup of old temp files (every hour)
def _cleanup_loop():
    while True:
        time.sleep(3600)
        try:
            cleanup_old_tasks()
        except Exception:
            pass

threading.Thread(target=_cleanup_loop, daemon=True).start()

if __name__ == '__main__':
    # Local testing only (Render uses Gunicorn)
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port)
